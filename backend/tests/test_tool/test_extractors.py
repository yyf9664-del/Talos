"""Unit tests for document content extractors."""

import pytest

from app.tool.extractors import extract_document, is_supported_binary


class TestIsSupportedBinary:
    def test_pdf(self):
        assert is_supported_binary("report.pdf")

    def test_docx(self):
        assert is_supported_binary("document.docx")

    def test_xlsx(self):
        assert is_supported_binary("data.xlsx")

    def test_pptx(self):
        assert is_supported_binary("deck.pptx")

    def test_case_insensitive(self):
        assert is_supported_binary("FILE.PDF")
        assert is_supported_binary("Doc.DOCX")

    def test_not_supported(self):
        assert not is_supported_binary("file.txt")
        assert not is_supported_binary("image.png")
        assert not is_supported_binary("code.py")
        assert not is_supported_binary("no_extension")

    def test_unsupported_extract_raises(self):
        with pytest.raises(ValueError, match="Unsupported format"):
            extract_document("file.txt")


class TestExtractPDF:
    @pytest.fixture(autouse=True)
    def _check_pypdf(self):
        pytest.importorskip("pypdf")

    def test_extract_pdf(self, tmp_path):
        from pypdf import PdfWriter
        from reportlab.pdfgen import canvas as rc

        # Create a simple PDF with text using reportlab
        pdf_path = tmp_path / "test.pdf"
        c = rc.Canvas(str(pdf_path))
        c.drawString(72, 700, "Hello PDF World")
        c.save()

        result = extract_document(str(pdf_path))
        assert "Hello PDF World" in result
        assert "Page 1" in result

    def test_extract_empty_pdf(self, tmp_path):
        from pypdf import PdfWriter

        pdf_path = tmp_path / "empty.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with open(pdf_path, "wb") as f:
            writer.write(f)

        result = extract_document(str(pdf_path))
        assert "No text content" in result


class TestExtractDOCX:
    @pytest.fixture(autouse=True)
    def _check_docx(self):
        pytest.importorskip("docx")

    def test_extract_docx(self, tmp_path):
        from docx import Document

        docx_path = tmp_path / "test.docx"
        doc = Document()
        doc.add_heading("Test Title", level=1)
        doc.add_paragraph("Hello DOCX World")
        doc.save(str(docx_path))

        result = extract_document(str(docx_path))
        assert "# Test Title" in result
        assert "Hello DOCX World" in result

    def test_extract_docx_with_table(self, tmp_path):
        from docx import Document

        docx_path = tmp_path / "table.docx"
        doc = Document()
        doc.add_paragraph("Before table")
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "A"
        table.rows[0].cells[1].text = "B"
        table.rows[1].cells[0].text = "1"
        table.rows[1].cells[1].text = "2"
        doc.save(str(docx_path))

        result = extract_document(str(docx_path))
        assert "Before table" in result
        assert "A | B" in result
        assert "1 | 2" in result

    def test_extract_empty_docx(self, tmp_path):
        from docx import Document

        docx_path = tmp_path / "empty.docx"
        doc = Document()
        doc.save(str(docx_path))

        result = extract_document(str(docx_path))
        assert "Empty document" in result


class TestExtractXLSX:
    @pytest.fixture(autouse=True)
    def _check_openpyxl(self):
        pytest.importorskip("openpyxl")

    def test_extract_xlsx(self, tmp_path):
        from openpyxl import Workbook

        xlsx_path = tmp_path / "test.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Alice"
        ws["B2"] = 42
        wb.save(str(xlsx_path))

        result = extract_document(str(xlsx_path))
        assert "Sheet: Data" in result
        assert "Name" in result
        assert "Alice" in result
        assert "42" in result

    def test_extract_xlsx_multi_sheet(self, tmp_path):
        from openpyxl import Workbook

        xlsx_path = tmp_path / "multi.xlsx"
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1["A1"] = "First"
        ws2 = wb.create_sheet("Sheet2")
        ws2["A1"] = "Second"
        wb.save(str(xlsx_path))

        result = extract_document(str(xlsx_path))
        assert "Sheet: Sheet1" in result
        assert "Sheet: Sheet2" in result
        assert "First" in result
        assert "Second" in result

    def test_extract_empty_xlsx(self, tmp_path):
        from openpyxl import Workbook

        xlsx_path = tmp_path / "empty.xlsx"
        wb = Workbook()
        wb.save(str(xlsx_path))

        result = extract_document(str(xlsx_path))
        assert "Empty workbook" in result


class TestExtractPPTX:
    @pytest.fixture(autouse=True)
    def _check_pptx(self):
        pytest.importorskip("pptx")

    def test_extract_pptx(self, tmp_path):
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        slide.placeholders[1].text = "Hello PPTX World"
        prs.save(str(pptx_path))

        result = extract_document(str(pptx_path))
        assert "Slide 1" in result
        assert "Test Slide" in result
        assert "Hello PPTX World" in result

    def test_extract_empty_pptx(self, tmp_path):
        from pptx import Presentation

        pptx_path = tmp_path / "empty.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        result = extract_document(str(pptx_path))
        assert "Empty presentation" in result
