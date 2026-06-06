"""Document content extractors for binary file formats.

Each extractor function takes a file path and returns extracted text.
Raises ImportError if the required library is not installed,
or other exceptions for corrupt/unreadable files.
"""

from __future__ import annotations

import os

# Extension -> extractor function name
_EXTRACTORS: dict[str, str] = {
    ".pdf": "_extract_pdf",
    ".docx": "_extract_docx",
    ".xlsx": "_extract_xlsx",
    ".pptx": "_extract_pptx",
}


def is_supported_binary(file_path: str) -> bool:
    """Return True if file has a supported binary extension."""
    ext = os.path.splitext(file_path)[1].lower()
    return ext in _EXTRACTORS


def extract_document(file_path: str) -> str:
    """Extract text from a supported binary document.

    Raises:
        ValueError: unsupported extension
        ImportError: required library not installed
        Exception: file corrupt or unreadable
    """
    ext = os.path.splitext(file_path)[1].lower()
    func_name = _EXTRACTORS.get(ext)
    if func_name is None:
        raise ValueError(f"Unsupported format: {ext}")
    func = globals()[func_name]
    return func(file_path)


# ------------------------------------------------------------------
# PDF
# ------------------------------------------------------------------

def _extract_pdf(file_path: str) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ImportError(
            "Cannot read PDF files: pypdf is not installed. "
            "Install with: pip install pypdf"
        )

    reader = PdfReader(file_path)
    pages: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"--- Page {i} ---\n{text}")

    if not pages:
        return "(No text content found in PDF — may be image-based or scanned)"

    return "\n\n".join(pages)


# ------------------------------------------------------------------
# DOCX
# ------------------------------------------------------------------

def _extract_docx(file_path: str) -> str:
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "Cannot read DOCX files: python-docx is not installed. "
            "Install with: pip install python-docx"
        )

    doc = Document(file_path)
    parts: list[str] = []

    # Paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = para.style.name if para.style else ""
        if style_name.startswith("Heading"):
            level_str = style_name.replace("Heading", "").strip()
            try:
                level = int(level_str)
            except ValueError:
                level = 1
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)

    # Tables
    for i, table in enumerate(doc.tables, 1):
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append(f"\n[Table {i}]")
            parts.append(rows[0])
            if len(rows) > 1:
                parts.append(" | ".join(["---"] * len(table.columns)))
                parts.extend(rows[1:])

    return "\n\n".join(parts) if parts else "(Empty document)"


# ------------------------------------------------------------------
# XLSX
# ------------------------------------------------------------------

def _extract_xlsx(file_path: str) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError(
            "Cannot read XLSX files: openpyxl is not installed. "
            "Install with: pip install openpyxl"
        )

    wb = load_workbook(file_path, read_only=True, data_only=True)
    sheets: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c for c in cells):  # skip entirely empty rows
                rows.append("\t".join(cells))
        if rows:
            sheets.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))

    wb.close()
    return "\n\n".join(sheets) if sheets else "(Empty workbook)"


# ------------------------------------------------------------------
# PPTX
# ------------------------------------------------------------------

def _extract_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "Cannot read PPTX files: python-pptx is not installed. "
            "Install with: pip install python-pptx"
        )

    prs = Presentation(file_path)
    slides: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        if texts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))

    return "\n\n".join(slides) if slides else "(Empty presentation)"
